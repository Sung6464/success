"""Extraction: turn an uploaded PDF, DOCX, or PPTX into a structured on-disk layout.

extracted_papers/<doc_id>/
    <doc_id>.md          # markdown w/ headings + ![](figures/...) anchors
    figures/             # extracted images
"""
from __future__ import annotations

import re
import shutil
from collections import Counter
from pathlib import Path

import fitz  # PyMuPDF

from . import config

SUPPORTED_IMAGE_EXTS = (
    "png", "jpg", "jpeg", "gif", "webp", "svg", "tiff", "tif", "bmp", "ico",
    "heic", "heif", "avif", "wmp"
)


def _safe_doc_id(filename: str) -> str:
    stem = Path(filename).stem
    stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem).strip("_")
    return stem or "document"


def _new_doc_dir(doc_id: str) -> Path:
    """Create a fresh <doc_id>/ folder (overwriting any previous extraction)."""
    out = config.PAPERS_DIR / doc_id
    if out.exists():
        shutil.rmtree(out)
    (out / "figures").mkdir(parents=True, exist_ok=True)
    (out / "tables").mkdir(parents=True, exist_ok=True)
    return out


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def _pdf_body_size(doc: "fitz.Document") -> float:
    """Most common font size = body text. Larger spans are treated as headings."""
    sizes: Counter = Counter()
    for page in doc:
        d = page.get_text("dict")
        for block in d.get("blocks", []):
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    txt = span.get("text", "").strip()
                    if txt:
                        sizes[round(span["size"], 1)] += len(txt)
    if not sizes:
        return 11.0
    return sizes.most_common(1)[0][0]


def _heading_level(size: float, body: float) -> int:
    """Map a font size to a markdown heading level (1..3), or 0 for body text."""
    ratio = size / body if body else 1.0
    if ratio >= 1.5:
        return 1
    if ratio >= 1.28:
        return 2
    if ratio >= 1.12:
        return 3
    return 0


def extract_pdf_with_document_intelligence(pdf_path: Path, doc_id: str) -> Path:
    """Extract text/tables/structure with Azure Document Intelligence, crop figures with PyMuPDF."""
    import base64
    from azure.ai.documentintelligence import DocumentIntelligenceClient
    from azure.core.credentials import AzureKeyCredential
    from azure.ai.documentintelligence.models import AnalyzeDocumentRequest

    out = _new_doc_dir(doc_id)
    fig_dir = out / "figures"

    client = DocumentIntelligenceClient(
        endpoint=config.DOCUMENT_INTELLIGENCE_ENDPOINT,
        credential=AzureKeyCredential(config.DOCUMENT_INTELLIGENCE_KEY)
    )

    pdf_bytes = pdf_path.read_bytes()
    try:
        poller = client.begin_analyze_document(
            model_id="prebuilt-layout",
            body=AnalyzeDocumentRequest(bytes_source=pdf_bytes),
            output_content_format="markdown"
        )
    except Exception:
        encoded_bytes = base64.b64encode(pdf_bytes).decode("utf-8")
        poller = client.begin_analyze_document(
            model_id="prebuilt-layout",
            body=AnalyzeDocumentRequest(base64_source=encoded_bytes),
            output_content_format="markdown"
        )

    result = poller.result()
    md_content = result.content or ""

    # Crop detected figures
    if result.figures:
        doc = fitz.open(pdf_path)
        for fig in result.figures:
            if fig.id and fig.bounding_regions:
                region = fig.bounding_regions[0]
                page_number = region.page_number
                polygon = region.polygon

                if polygon and len(polygon) == 8 and 0 < page_number <= len(doc):
                    page = doc[page_number - 1]
                    x_coords = [polygon[i] for i in range(0, 8, 2)]
                    y_coords = [polygon[i+1] for i in range(0, 8, 2)]

                    min_x = min(x_coords) * 72.0
                    min_y = min(y_coords) * 72.0
                    max_x = max(x_coords) * 72.0
                    max_y = max(y_coords) * 72.0

                    rect = fitz.Rect(min_x, min_y, max_x, max_y)
                    try:
                        pix = page.get_pixmap(clip=rect, dpi=150)
                        fname = f"img_{fig.id}.png"
                        pix.save(str(fig_dir / fname))
                    except Exception as e:
                        print(f"Failed to crop figure {fig.id}: {e}")
        doc.close()

    # Rewrite Markdown image URLs
    md_content = re.sub(r"\!\[(.*?)\]\(figures/([^)]+)\)", r"![\1](figures/img_\2.png)", md_content)

    if not md_content.lstrip().startswith("#"):
        md_content = f"# {doc_id}\n\n" + md_content

    (out / f"{doc_id}.md").write_text(md_content, encoding="utf-8")
    return out


def extract_pdf(pdf_path: Path, doc_id: str) -> Path:
    if config.has_document_intelligence():
        try:
            return extract_pdf_with_document_intelligence(pdf_path, doc_id)
        except Exception as e:
            print(f"Azure Document Intelligence failed: {e}. Falling back to local PyMuPDF extraction...")

    out = _new_doc_dir(doc_id)
    fig_dir = out / "figures"
    doc = fitz.open(pdf_path)
    body = _pdf_body_size(doc)

    md_lines: list[str] = []
    img_counter = 0
    seen_xrefs: dict[int, str] = {}

    for page in doc:
        page_dict = page.get_text("dict")
        items: list[tuple[float, str]] = []

        for block in page_dict.get("blocks", []):
            btype = block.get("type", 0)
            y0 = block.get("bbox", [0, 0, 0, 0])[1]

            if btype == 0:  # text
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    text = "".join(s.get("text", "") for s in spans).strip()
                    if not text:
                        continue
                    max_size = max((s.get("size", body) for s in spans), default=body)
                    is_bold = any("bold" in (s.get("font", "").lower()) for s in spans)
                    lvl = _heading_level(max_size, body)
                    if lvl == 0 and is_bold and len(text) < 80 and not text.endswith("."):
                        lvl = 3
                    if lvl:
                        items.append((y0, f"{'#' * lvl} {text}"))
                    else:
                        items.append((y0, text))

            elif btype == 1:  # image
                xref = block.get("number")
                img_bytes = block.get("image")
                ext = block.get("ext", "png")
                if not img_bytes:
                    continue
                key = (xref if isinstance(xref, int) else id(img_bytes))
                if key in seen_xrefs:
                    rel = seen_xrefs[key]
                else:
                    img_counter += 1
                    ext_to_use = ext.lower() if ext.lower() in SUPPORTED_IMAGE_EXTS else "png"
                    fname = f"img_{img_counter}.{ext_to_use}"
                    (fig_dir / fname).write_bytes(img_bytes)
                    rel = f"figures/{fname}"
                    seen_xrefs[key] = rel
                items.append((y0, f"![]({rel})"))

        # Fallback: some PDFs expose images only via get_images()
        if not any(frag.startswith("![]") for _, frag in items):
            for info in page.get_images(full=True):
                xref = info[0]
                if xref in seen_xrefs:
                    continue
                try:
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n - pix.alpha >= 4:  # CMYK -> RGB
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_counter += 1
                    fname = f"img_{img_counter}.png"
                    pix.save(str(fig_dir / fname))
                    rel = f"figures/{fname}"
                    seen_xrefs[xref] = rel
                    items.append((9_999, f"![]({rel})"))
                    pix = None
                except Exception:
                    continue

        items.sort(key=lambda t: t[0])
        for _, frag in items:
            md_lines.append(frag)
            md_lines.append("")

    doc.close()
    md = _tidy_headings("\n".join(md_lines))
    if not md.lstrip().startswith("#"):
        md = f"# {doc_id}\n\n" + md
    (out / f"{doc_id}.md").write_text(md, encoding="utf-8")
    return out


def _tidy_headings(md: str) -> str:
    lines = [ln for ln in md.split("\n")]
    out: list[str] = []
    i = 0
    head_re = re.compile(r"^(#{1,6})\s+(.*)$")
    num_re = re.compile(r"^[\d]+(\.[\d]+)*\.?$")
    while i < len(lines):
        ln = lines[i].rstrip()
        m = head_re.match(ln)
        if m:
            hashes, txt = m.group(1), m.group(2).strip()
            if len(txt) <= 1 and not txt.isalnum():
                i += 1
                continue
            if num_re.match(txt):
                j = i + 1
                while j < len(lines) and not lines[j].strip():
                    j += 1
                if j < len(lines):
                    m2 = head_re.match(lines[j].rstrip())
                    if m2:
                        merged = f"{m2.group(1)} {txt} {m2.group(2).strip()}"
                        out.append(merged)
                        out.append("")
                        i = j + 1
                        continue
            out.append(f"{hashes} {txt}")
        else:
            out.append(ln)
        i += 1
    return "\n".join(out).strip() + "\n"


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------
def extract_docx(docx_path: Path, doc_id: str) -> Path:
    import docx
    from docx.oxml.ns import qn

    out = _new_doc_dir(doc_id)
    fig_dir = out / "figures"
    document = docx.Document(str(docx_path))

    md_lines: list[str] = []
    img_counter = 0
    rels = document.part.rels

    def save_blip(rid: str) -> str | None:
        nonlocal img_counter
        try:
            part = rels[rid].target_part
        except Exception:
            return None
        blob = part.blob
        ext = (part.partname.ext or ".png").lstrip(".")
        if ext.lower() not in SUPPORTED_IMAGE_EXTS:
            ext = "png"
        img_counter += 1
        fname = f"img_{img_counter}.{ext.lower()}"
        (fig_dir / fname).write_bytes(blob)
        return f"figures/{fname}"

    for para in document.paragraphs:
        text = para.text.strip()
        style = (para.style.name or "").lower() if para.style else ""

        para_imgs: list[str] = []
        for blip in para._p.findall(".//" + qn("a:blip")):
            rid = blip.get(qn("r:embed")) or blip.get(qn("r:link"))
            if rid:
                rel = save_blip(rid)
                if rel:
                    para_imgs.append(rel)

        if text:
            if style.startswith("title"):
                md_lines += [f"# {text}", ""]
            elif style.startswith("heading"):
                m = re.search(r"(\d+)", style)
                lvl = min(int(m.group(1)), 6) if m else 2
                md_lines += [f"{'#' * lvl} {text}", ""]
            else:
                md_lines += [text, ""]

        for rel in para_imgs:
            md_lines += [f"![]({rel})", ""]

    for tbl in document.tables:
        rows = []
        for row in tbl.rows:
            rows.append([c.text.strip().replace("\n", " ") for c in row.cells])
        if not rows:
            continue
        header = rows[0]
        md_lines.append("| " + " | ".join(header) + " |")
        md_lines.append("| " + " | ".join("---" for _ in header) + " |")
        for r in rows[1:]:
            md_lines.append("| " + " | ".join(r) + " |")
        md_lines.append("")

    md = _tidy_headings("\n".join(md_lines))
    if not md.lstrip().startswith("#"):
        md = f"# {doc_id}\n\n" + md
    (out / f"{doc_id}.md").write_text(md, encoding="utf-8")
    return out


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
def extract_pptx(pptx_path: Path, doc_id: str) -> Path:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out = _new_doc_dir(doc_id)
    fig_dir = out / "figures"
    prs = Presentation(str(pptx_path))

    md_lines: list[str] = []
    img_counter = 0

    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        if not title:
            title = f"Slide {i + 1}"

        md_lines.append(f"## {title}")
        md_lines.append("")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape == slide.shapes.title:
                    continue
                md_lines.append(shape.text.strip())
                md_lines.append("")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_counter += 1
                image = shape.image
                ext = image.ext
                fname = f"img_{img_counter}.{ext}"
                (fig_dir / fname).write_bytes(image.blob)
                md_lines.append(f"![](figures/{fname})")
                md_lines.append("")

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text_frame.text.strip().replace("\n", " ") for cell in row.cells])
                if rows:
                    header = rows[0]
                    md_lines.append("| " + " | ".join(header) + " |")
                    md_lines.append("| " + " | ".join("---" for _ in header) + " |")
                    for r in rows[1:]:
                        md_lines.append("| " + " | ".join(r) + " |")
                    md_lines.append("")

    md = "\n".join(md_lines)
    md = f"# {doc_id}\n\n" + md

    (out / f"{doc_id}.md").write_text(md, encoding="utf-8")
    return out


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------
def extract_file(path: str | Path, original_name: str | None = None) -> tuple[str, Path]:
    path = Path(path)
    name = original_name or path.name
    doc_id = _safe_doc_id(name)
    suffix = Path(name).suffix.lower()

    if suffix == ".pdf":
        out = extract_pdf(path, doc_id)
    elif suffix in (".docx", ".doc"):
        out = extract_docx(path, doc_id)
    elif suffix in (".pptx", ".ppt"):
        out = extract_pptx(path, doc_id)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")
    return doc_id, out
