import os
import io
import logging
from typing import List, Optional
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """
    Handles document loading, processing, and chunking for RAG system.
    
    Strategy:
    - Each file = One complete story
    - Chunk stories for better matching
    - Preserve story_id in metadata
    - When chunks match, return complete story
    """
    
    def __init__(self, blob_manager, chunk_size: int = 1000, chunk_overlap: int = 200, enable_ocr: bool = True):
        """
        Initialize Document Processor.
        
        Args:
            blob_manager: BlobStorageManager instance
            chunk_size: Size of each text chunk (default: 1000)
            chunk_overlap: Overlap between chunks (default: 200)
            enable_ocr: Enable OCR for extracting text from images (default: False)
        """
        self.blob_manager = blob_manager
        self.enable_ocr = enable_ocr
        
        # Chunking configuration for optimal retrieval
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        # Check OCR availability
        if self.enable_ocr:
            try:
                import pytesseract
                from PIL import Image
                logger.info("✅ OCR enabled (Tesseract available)")
            except ImportError:
                logger.warning("⚠️ OCR requested but dependencies not installed. Install: pip install pytesseract pillow pdf2image")
                self.enable_ocr = False
        
        logger.info(f"DocumentProcessor initialized with chunk_size={chunk_size}, chunk_overlap={chunk_overlap}, OCR={'enabled' if self.enable_ocr else 'disabled'}")
    
    def load_document_from_blob(self, blob_name: str) -> Optional[Document]:
        """
        Load a document from blob storage and extract text.
        
        Args:
            blob_name: Name of the blob file
            
        Returns:
            Document object with text and metadata, or None if failed
        """
        try:
            # Download file content from blob
            file_content = self.blob_manager.download_blob_to_memory(blob_name)
            file_extension = os.path.splitext(blob_name)[1].lower()
            
            # Extract text based on file type
            if file_extension == '.pdf':
                text = self._extract_pdf_text(file_content)
            elif file_extension in ['.docx', '.doc']:
                text = self._extract_docx_text(file_content)
            elif file_extension in ['.pptx', '.ppt']:
                text = self._extract_pptx_text(file_content)
            elif file_extension == '.txt':
                text = file_content.decode('utf-8')
            else:
                logger.warning(f"Unsupported file type: {file_extension} for {blob_name}")
                return None
            
            if not text or not text.strip():
                logger.warning(f"No text extracted from {blob_name}")
                return None
            
            # Create document with metadata
            doc = Document(
                page_content=text,
                metadata={
                    'source': blob_name,
                    'source_file': blob_name,  # Add source_file for consistency
                    'story_id': self._get_story_id(blob_name),
                    'category': self._get_category(blob_name),
                    'file_type': file_extension,
                    'blob_url': self.blob_manager.get_blob_url(blob_name)
                }
            )
            
            logger.info(f"✅ Loaded document: {blob_name} ({len(text)} characters)")
            return doc
            
        except Exception as e:
            logger.error(f"❌ Error loading document {blob_name}: {e}")
            return None
    
    def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF file content, including OCR if enabled."""
        try:
            from pypdf import PdfReader
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)
            text = ""
            
            # Extract text from each page
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                
                # If no text found and OCR is enabled, try OCR
                if self.enable_ocr and (not page_text or len(page_text.strip()) < 50):
                    logger.info(f"  Page {page_num + 1}: Low text content, attempting OCR...")
                    ocr_text = self._ocr_pdf_page(file_content, page_num)
                    if ocr_text:
                        text += f"\n[OCR Page {page_num + 1}]\n{ocr_text}\n"
                    else:
                        text += page_text + "\n"
                else:
                    text += page_text + "\n"
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def _ocr_pdf_page(self, file_content: bytes, page_num: int) -> str:
        """Extract text from PDF page using OCR."""
        try:
            from pdf2image import convert_from_bytes
            import pytesseract
            
            # Convert specific page to image
            images = convert_from_bytes(file_content, first_page=page_num + 1, last_page=page_num + 1)
            
            if not images:
                return ""
            
            # Perform OCR on the image
            text = pytesseract.image_to_string(images[0])
            return text.strip()
        except Exception as e:
            logger.error(f"Error performing OCR on PDF page {page_num}: {e}")
            return ""
    
    def _extract_docx_text(self, file_content: bytes) -> str:
        """Extract text from DOCX file content, including images if OCR enabled."""
        try:
            from docx import Document as DocxDocument
            docx_file = io.BytesIO(file_content)
            doc = DocxDocument(docx_file)
            
            # Extract text from paragraphs
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Extract text from images if OCR enabled
            if self.enable_ocr:
                image_text = self._ocr_docx_images(doc)
                if image_text:
                    text += "\n\n[Extracted from Images]\n" + image_text
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    def _ocr_docx_images(self, doc) -> str:
        """Extract text from images in DOCX using OCR."""
        try:
            import pytesseract
            from PIL import Image
            
            image_text = ""
            image_count = 0
            
            # Iterate through document parts to find images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        image = Image.open(io.BytesIO(image_data))
                        
                        # Perform OCR
                        text = pytesseract.image_to_string(image)
                        if text.strip():
                            image_count += 1
                            image_text += f"\n[Image {image_count}]\n{text}\n"
                    except Exception as e:
                        logger.debug(f"Could not OCR image: {e}")
            
            if image_count > 0:
                logger.info(f"  Extracted text from {image_count} images in DOCX")
            
            return image_text.strip()
        except Exception as e:
            logger.error(f"Error performing OCR on DOCX images: {e}")
            return ""
    
    def _extract_pptx_text(self, file_content: bytes) -> str:
        """Extract text from PPTX file content, including images if OCR enabled."""
        try:
            from pptx import Presentation
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            text = ""
            
            # Iterate through slides
            for slide_num, slide in enumerate(prs.slides, 1):
                # Add slide separator
                text += f"\n--- Slide {slide_num} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text += cell.text + " "
                            text += "\n"
                    
                    # Extract text from text frames
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text:
                                    text += run.text
                            text += "\n"
                    
                    # Extract text from images if OCR enabled
                    if self.enable_ocr and shape.shape_type == 13:  # 13 = Picture
                        try:
                            image_text = self._ocr_pptx_image(shape)
                            if image_text:
                                text += f"[Image Text]\n{image_text}\n"
                        except Exception as e:
                            logger.debug(f"Could not OCR image in slide {slide_num}: {e}")
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _ocr_pptx_image(self, shape) -> str:
        """Extract text from a PowerPoint image using OCR."""
        try:
            import pytesseract
            from PIL import Image
            
            # Get image data
            image_blob = shape.image.blob
            image = Image.open(io.BytesIO(image_blob))
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.debug(f"Error performing OCR on PPTX image: {e}")
            return ""
    
    def _get_story_id(self, blob_name: str) -> str:
        """Generate a unique story ID from blob name."""
        # Remove file extension and path separators
        story_id = os.path.splitext(blob_name)[0].replace('/', '_').replace('\\', '_')
        return story_id
    
    def _get_category(self, blob_name: str) -> str:
        """
        Extract category from blob path.
        Assumes folder structure like: category/filename.ext
        
        Args:
            blob_name: Full blob path (e.g., 'insurance/story.pdf')
            
        Returns:
            Category name (e.g., 'insurance') or 'general' if no folder
        """
        # Split by forward slash or backslash
        parts = blob_name.replace('\\', '/').split('/')
        
        # If there's a folder before the filename, use it as category
        if len(parts) >= 2:
            folder_name = parts[-2]
        else:
            folder_name = "general"
        return folder_name
    
    def chunk_document(self, document: Document) -> List[Document]:
        """
        Split document into chunks while preserving metadata.
        
        Each chunk will have:
        - Original metadata (story_id, source, etc.)
        - Chunk index for ordering
        
        Args:
            document: Document to chunk
            
        Returns:
            List of chunked Documents
        """
        chunks = self.text_splitter.split_documents([document])
        
        # Add chunk index to metadata
        for idx, chunk in enumerate(chunks):
            chunk.metadata['chunk_index'] = idx
            chunk.metadata['total_chunks'] = len(chunks)
        
        logger.info(f"  Created {len(chunks)} chunks from {document.metadata.get('source', 'unknown')}")
        return chunks
    
    def process_all_blobs(self, prefix: Optional[str] = None) -> dict:
        """
        Process all blobs from storage and return documents with chunks.
        
        Args:
            prefix: Optional prefix to filter blobs (e.g., 'success_stories/')
            
        Returns:
            Dictionary with 'documents' and 'chunks' lists
        """
        # Get all blobs
        blobs = self.blob_manager.list_all_blobs(prefix=prefix)
        logger.info(f"📁 Found {len(blobs)} files to process")
        
        all_documents = []
        all_chunks = []
        
        for blob in blobs:
            blob_name = blob['name']
            
            # Load document
            doc = self.load_document_from_blob(blob_name)
            
            if not doc:
                continue
            
            all_documents.append(doc)
            
            # Chunk document
            chunks = self.chunk_document(doc)
            all_chunks.extend(chunks)
        
        logger.info(f"✅ Processed {len(all_documents)} documents into {len(all_chunks)} chunks")
        
        return {
            'documents': all_documents,
            'chunks': all_chunks
        }


# if __name__ == "__main__":
#     # Test the document processor
#     from blob_manager_v1 import BlobStorageManager
    
#     print("Testing DocumentProcessor...")
    
#     blob_manager = BlobStorageManager()
#     processor = DocumentProcessor(blob_manager)
    
#     # List blobs
#     blobs = blob_manager.list_all_blobs()
#     if blobs:
#         print(f"\nTesting with first file: {blobs[0]['name']}")
#         doc = processor.load_document_from_blob(blobs[0]['name'])
#         if doc:
#             chunks = processor.chunk_document(doc)
#             print(f"Created {len(chunks)} chunks")
#             print(f"First chunk: {chunks[0].page_content[:200]}...")
